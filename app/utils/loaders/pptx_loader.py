from pathlib import Path
from typing import List

from pptx import Presentation
from langchain_core.documents import Document

from app.utils.base.document_loader import DocumentLoader


class PPTXLoader(DocumentLoader):
    def load(self, file_path: Path) -> List[Document]:
        def _load():
            text = f"File: {file_path.name}\n\n"
            try:
                prs = Presentation(str(file_path))
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = f"Slide {i}:\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text += shape.text.strip() + "\n"
                    text += slide_text + "\n"
            except Exception as e:
                text += f"[Lỗi đọc PPTX: {str(e)}]\n"

            metadata = self._create_metadata(file_path, "pptx")
            return [Document(page_content=text, metadata=metadata)]
        return self._safe_load(file_path, _load)